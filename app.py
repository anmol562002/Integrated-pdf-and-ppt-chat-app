import streamlit as st
import base64
import openai
import pptx
from pptx.util import Inches, Pt
import os
import pickle
from PyPDF2 import PdfReader
from streamlit_extras.add_vertical_space import add_vertical_space
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.llms import OpenAI
from langchain.chains.question_answering import load_qa_chain

from dotenv import load_dotenv
load_dotenv()

# Set OpenAI API key
openai.api_key = os.getenv('OPENAI_API_KEY')  # Replace with your actual API key

# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

# Define custom formatting options for PDF Chat
PDF_CHAT_TITLE = "Chat with babaji💬"

# Sidebar contents
with st.sidebar:
    st.title('🤗💬 babaji App by Boyanmol')
    st.markdown('''
    ## About
    This app is powered by:
    - [Streamlit](https://streamlit.io/)
    - [LangChain](https://python.langchain.com/)
    - [OpenAI](https://platform.openai.com/docs/models) LLM model
 
    ''')
    add_vertical_space(5)
    st.write('Made with ❤️ by [boyanmol](https://youtube.com/@boyaanmol)')

def generate_slide_titles(topic):
    prompt = f"Generate 5 slide titles for the topic '{topic}'."
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=300,
    )
    return response['choices'][0]['text'].split("\n")

def generate_slide_content(slide_title):
    prompt = f"Generate content for the slide: '{slide_title}'."
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=300,  # Adjust as needed based on the desired content length
    )
    return response['choices'][0]['text']

def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        # Customize font size for titles and content
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    ppt_filename = f"generated_ppt/{topic}_presentation.pptx"
    prs.save(ppt_filename)
    return ppt_filename

def main():
    st.title("AI PDF and PPT webapp ")

    # Create a sidebar for navigation
    app_mode = st.sidebar.selectbox("Select App Mode", ("PPT Generator", "Chat with Pdf"))

    if app_mode == "PPT Generator":
        st.header("PPT Generator")

        topic = st.text_input("Enter the topic for your presentation:")
        generate_button = st.button("Generate Presentation")

        if generate_button and topic:
            st.info("Generating presentation... Please wait.")
            slide_titles = generate_slide_titles(topic)
            filtered_slide_titles = [item for item in slide_titles if item.strip() != '']
            slide_contents = [generate_slide_content(title) for title in filtered_slide_titles]
            ppt_filename = create_presentation(topic, filtered_slide_titles, slide_contents)

            st.success("Presentation generated successfully!")
            st.markdown(get_ppt_download_link(ppt_filename), unsafe_allow_html=True)

    elif app_mode == "Chat with Pdf":
        st.header(PDF_CHAT_TITLE)

        load_dotenv()

         # upload a PDF file
        pdf = st.file_uploader("Upload your PDF", type='pdf')

        if pdf is not None:
            pdf_reader = PdfReader(pdf)
            
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()

            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len
                )
            chunks = text_splitter.split_text(text=text)

            # # embeddings
            
            store_name = pdf.name[:-4]

            if os.path.exists(f"{store_name}.pkl"):
                with open(f"{store_name}.pkl", "rb") as f:
                    VectorStore = pickle.load(f)
                #st.write('Embeddings Loaded from the Disk')
            else:
                embeddings = OpenAIEmbeddings()
                VectorStore = FAISS.from_texts(chunks, embedding=embeddings)
                with open(f"{store_name}.pkl", "wb") as f:
                    pickle.dump(VectorStore, f)
                    #st.write('Embeddings computation completed')
            
            # Accept user questions/query
            query = st.text_input("Ask questions about your PDF file so panda tell truth it don't lie :")
            #st.write(query)

            if query:
                docs = VectorStore.similarity_search(query=query, k=3)
            
                llm = OpenAI(temperature=0,)
                chain = load_qa_chain(llm=llm, chain_type="stuff")
                response = chain.run(input_documents=docs, question=query)
                st.write(response)

def get_ppt_download_link(ppt_filename):
    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'

if __name__ == "__main__":
    main()